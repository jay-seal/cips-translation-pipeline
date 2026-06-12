"""
cips_reconstruct.py
====================
CIPS Translation Pipeline — PPTX Reconstruction Script

Triggered by GitHub Actions (reconstruct.yml). Accepts file paths via
command-line arguments. All file I/O is handled by the workflow.

Usage:
    python cips_reconstruct.py \
        --source  inputs/source.pptx \
        --translations  inputs/agent3_output.json \
        [--qa  inputs/agent4_output.json] \
        --output  outputs/<filename>.pptx \
        [--failure-threshold 0.30]

Outputs written to the same directory as --output:
    <name>.pptx                     Translated presentation.
    <name>_match_report.json        Full per-segment match results (machine-readable).
    <name>_manual_corrections.html  Human-readable list of unmatched segments for
                                    manual correction in PowerPoint.

Matching strategy — applied in order per segment, stopping at first success:

    Tier 1   Exact normalised match on the full shape text.
    Tier 1b  Page-number-tolerant match. Used for footer shapes where the
             slide-number field element has an empty or mismatched cached
             value: strips trailing "| N" from both source and shape text
             before comparing. Each slide's footer shape is updated
             individually with its own correct static page number.
    Tier 2   Substring normalised match (single-paragraph shapes only).
    Tier 3   Per-paragraph match — replaces only the matched paragraph.
    Tier 3b  Contiguous paragraph range match — for multi-line segments,
             tries matching N consecutive paragraphs against the N lines of
             source_text. Handles quotation blocks, agenda lists, and other
             multi-paragraph segments extracted from a larger shape.
    Tier 4   Layout shapes — Tiers 1–1b–2–3–3b on slide.slide_layout.shapes.
    Tier 5   Master shapes — Tiers 1–1b–2–3–3b on slide_layout.slide_master.shapes.
    Notes    Speaker notes — Tiers 1–3–3b on slide.notes_slide.notes_text_frame.
             Applied only for segments with element_type='speaker_note'.

    Once a layout or master shape is translated, all subsequent slides that
    reference the same text return LAYOUT_ALREADY_TRANSLATED or
    MASTER_ALREADY_TRANSLATED rather than NO_MATCH.

Normalisation (comparison only — never applied to output text):
    1. NFKC unicode normalisation.
    2. Smart quote / typographic punctuation normalisation.
    3. Collapse all whitespace to a single ASCII space.
    4. Strip leading/trailing whitespace.
    5. Strip leading bullet/list characters (•, ❶–❿ etc.).

Field element handling:
    All field elements (<a:fld>) in replaced paragraphs are removed and the
    full translated text (including any literal page number) is written as
    static text.

Non-text-box segment filtering:
    Segments where is_in_text_box=False and whose element_type is NOT in
    _ALWAYS_PROCESS_TYPES are excluded from the matching loop and logged as
    SKIP_NOT_TEXT_BOX.

Speaker notes:
    Segments with element_type='speaker_note' are matched and replaced in
    slide.notes_slide.notes_text_frame. They bypass the main shape-matching
    loop. Speaker note NO_MATCH segments are included in the manual
    corrections HTML since they are accessible via PowerPoint's Notes view.

Known permanent failures:
    Text embedded in raster images or SmartArt is not accessible to
    python-pptx and will always produce NO_MATCH.

Exit codes:
    0   Success.
    1   Missing or invalid input files.
    2   Match failure rate exceeds --failure-threshold.
"""

import argparse
import html
import json
import logging
import re
import sys
import unicodedata
from datetime import datetime, timezone
from pathlib import Path

from lxml import etree
from pptx import Presentation

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("cips_reconstruct")

# ---------------------------------------------------------------------------
# XML namespace constants
# ---------------------------------------------------------------------------
_A_NS  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_A_T   = f'{{{_A_NS}}}t'
_A_R   = f'{{{_A_NS}}}r'
_A_FLD = f'{{{_A_NS}}}fld'

# ---------------------------------------------------------------------------
# Bullet/list character stripping
# ---------------------------------------------------------------------------
_LEADING_BULLET_RE = re.compile(
    '['
    '\u2022'        # BULLET •
    '\u2023'        # TRIANGULAR BULLET
    '\u25E6'        # WHITE BULLET
    '\u2043'        # HYPHEN BULLET
    '\u2219'        # BULLET OPERATOR
    '\u25CF'        # BLACK CIRCLE
    '\u25AA'        # BLACK SMALL SQUARE
    '\u2776-\u277F' # DINGBAT NEGATIVE CIRCLED DIGITS ❶-❿
    '\u2780-\u2789' # DINGBAT CIRCLED SANS-SERIF DIGITS
    ']+' + r'\s*'
)

_TRAILING_PAGE_RE = re.compile(r'\s*\|\s*\d+\s*$')

# ---------------------------------------------------------------------------
# Element types that are always editable PPTX content shapes.
# 'label' is included: labels appear in accessible TextBox/Rectangle/AutoShape
# elements as well as inaccessible SmartArt. Including it ensures accessible
# label shapes are translated; inaccessible ones produce NO_MATCH (visible in
# the corrections report) rather than being silently skipped.
# 'speaker_note' is NOT included: notes are handled via a dedicated path.
# ---------------------------------------------------------------------------
_ALWAYS_PROCESS_TYPES = frozenset({
    'slide_title',
    'slide_subtitle',
    'body_text',
    'bullet_point',
    'heading',
    'table_cell',
    'text_box',
    'caption',
    'label',
})


# ---------------------------------------------------------------------------
# Text normalisation
# ---------------------------------------------------------------------------

def normalise(text: str) -> str:
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    text = text.replace('\u2018', "'").replace('\u2019', "'")
    text = text.replace('\u201C', '"').replace('\u201D', '"')
    text = text.replace('\u2013', '-').replace('\u2014', '-')
    text = re.sub(r"\s+", " ", text)
    text = text.strip()
    text = _LEADING_BULLET_RE.sub('', text)
    return text


def _source_key(norm_source: str) -> str:
    return _TRAILING_PAGE_RE.sub('', norm_source).strip()


# ---------------------------------------------------------------------------
# Paragraph text helpers
# ---------------------------------------------------------------------------

def _para_full_text(para) -> str:
    return ''.join(elem.text or '' for elem in para._p.iter(_A_T))


def shape_full_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(_para_full_text(p) for p in shape.text_frame.paragraphs)


def paragraph_text(para) -> str:
    return _para_full_text(para)


# ---------------------------------------------------------------------------
# Table cell proxy
# ---------------------------------------------------------------------------

class _TableCellProxy:
    __slots__ = ('has_text_frame', 'text_frame', 'name')

    def __init__(self, cell, row_idx: int, col_idx: int, parent_name: str):
        self.has_text_frame = True
        self.text_frame = cell.text_frame
        self.name = f"{parent_name}[r{row_idx}c{col_idx}]"


# ---------------------------------------------------------------------------
# Notes proxy
# ---------------------------------------------------------------------------

class _NotesProxy:
    __slots__ = ('has_text_frame', 'text_frame', 'name')

    def __init__(self, text_frame, slide_id):
        self.has_text_frame = True
        self.text_frame = text_frame
        self.name = f"notes_slide_{slide_id}"


# ---------------------------------------------------------------------------
# Shape iteration
# ---------------------------------------------------------------------------

def _iter_shapes(shape_collection):
    for shape in shape_collection:
        st = shape.shape_type
        if st == 6:
            yield from _iter_group(shape)
        elif st == 19:
            yield from _iter_table(shape)
        else:
            yield shape


def _iter_group(group_shape):
    for shape in group_shape.shapes:
        st = shape.shape_type
        if st == 6:
            yield from _iter_group(shape)
        elif st == 19:
            yield from _iter_table(shape)
        else:
            yield shape


def _iter_table(table_shape):
    for ri, row in enumerate(table_shape.table.rows):
        for ci, cell in enumerate(row.cells):
            yield _TableCellProxy(cell, ri, ci, table_shape.name)


def _shape_label(shape) -> str:
    return getattr(shape, 'name', '?')


# ---------------------------------------------------------------------------
# Run-level text replacement
# ---------------------------------------------------------------------------

def _replace_runs_in_paragraph(paragraph, new_text: str) -> None:
    p_elem = paragraph._p
    runs = paragraph.runs
    field_elements = p_elem.findall(_A_FLD)

    if runs:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
        for fld in field_elements:
            p_elem.remove(fld)
    elif field_elements:
        r_elem = etree.SubElement(p_elem, _A_R)
        t_elem = etree.SubElement(r_elem, _A_T)
        t_elem.text = new_text
        for fld in field_elements:
            p_elem.remove(fld)


def replace_shape_text(shape, translated_text: str) -> None:
    if not shape.has_text_frame:
        return
    paragraphs = shape.text_frame.paragraphs
    if not paragraphs:
        return
    translated_lines = translated_text.split("\n")
    non_empty_paras = [p for p in paragraphs if _para_full_text(p).strip()]
    if len(non_empty_paras) > 1 and len(translated_lines) != len(non_empty_paras):
        log.warning(
            "LINE_COUNT_MISMATCH  shape=%-30s  "
            "shape_paragraphs=%d  translated_lines=%d",
            _shape_label(shape), len(non_empty_paras), len(translated_lines),
        )
    for i, para in enumerate(paragraphs):
        if i < len(translated_lines):
            _replace_runs_in_paragraph(para, translated_lines[i])
        else:
            _replace_runs_in_paragraph(para, "")


def replace_paragraph_text(shape, para_index: int, translated_text: str) -> None:
    if not shape.has_text_frame:
        return
    paragraphs = shape.text_frame.paragraphs
    if para_index >= len(paragraphs):
        return
    _replace_runs_in_paragraph(paragraphs[para_index], translated_text)


def replace_paragraph_range(shape, start: int, count: int,
                             translated_text: str) -> None:
    """Replace `count` consecutive paragraphs starting at `start`."""
    if not shape.has_text_frame:
        return
    paragraphs = shape.text_frame.paragraphs
    translated_lines = translated_text.split("\n")
    for i in range(count):
        idx = start + i
        if idx >= len(paragraphs):
            break
        line = translated_lines[i] if i < len(translated_lines) else ""
        _replace_runs_in_paragraph(paragraphs[idx], line)


# ---------------------------------------------------------------------------
# Single-shape matching helper
# ---------------------------------------------------------------------------

def _match_shape(shape, norm_source: str, translated_text: str,
                 segment_id: str, slide_id, tier_prefix: str = "",
                 source_text_raw: str = ""):
    """
    Apply Tiers 1, 1b, 2, 3, and 3b to a single shape.
    Returns a result dict on match, else None.

    Tier 3b — Contiguous paragraph range match:
    For multi-line segments (source_text contains newlines), tries to find N
    consecutive paragraphs in the shape whose normalised texts exactly match
    the N lines of source_text. This handles quotation blocks, agenda lists,
    and other multi-paragraph segments that were extracted from a larger shape
    and cannot be matched by T1 (full shape) or T3 (single paragraph).
    """
    if not shape.has_text_frame:
        return None

    norm_shape = normalise(shape_full_text(shape))
    label = _shape_label(shape)

    # Tier 1 — Exact normalised match.
    if norm_shape == norm_source:
        replace_shape_text(shape, translated_text)
        log.info("%-24s seg=%-12s  slide=%s  shape=%s",
                 f"{tier_prefix}T1_EXACT", segment_id, slide_id, label)
        return {"result": f"{tier_prefix}T1_EXACT", "shape": label}

    # Tier 1b — Page-number-tolerant match.
    src_key   = _source_key(norm_source)
    shape_key = _source_key(norm_shape)
    if (src_key
            and src_key != norm_source
            and src_key == shape_key):
        replace_shape_text(shape, translated_text)
        log.info("%-24s seg=%-12s  slide=%s  shape=%s",
                 f"{tier_prefix}T1b_FUZZY", segment_id, slide_id, label)
        return {"result": f"{tier_prefix}T1b_FUZZY", "shape": label}

    # Tier 2 — Substring match (single-paragraph shapes only).
    non_empty_paras = [
        p for p in shape.text_frame.paragraphs
        if normalise(paragraph_text(p))
    ]
    if (len(norm_source) > 3
            and norm_source in norm_shape
            and len(non_empty_paras) == 1):
        replace_shape_text(shape, translated_text)
        log.info("%-24s seg=%-12s  slide=%s  shape=%s",
                 f"{tier_prefix}T2_SUBSTR", segment_id, slide_id, label)
        return {"result": f"{tier_prefix}T2_SUBSTRING", "shape": label}

    # Tier 3 — Per-paragraph match (single-line source texts).
    for idx, para in enumerate(shape.text_frame.paragraphs):
        norm_para = normalise(paragraph_text(para))
        if norm_para == norm_source and len(norm_source) > 1:
            replace_paragraph_text(shape, idx, translated_text)
            log.info("%-24s seg=%-12s  slide=%s  shape=%s  para=%d",
                     f"{tier_prefix}T3_PARA", segment_id, slide_id, label, idx)
            return {"result": f"{tier_prefix}T3_PARAGRAPH", "shape": label,
                    "para_index": idx}

    # Tier 3b — Contiguous paragraph range match (multi-line source texts).
    # IMPORTANT: norm_source has had all whitespace (including \n) collapsed to
    # spaces, so we cannot split norm_source to recover the original lines.
    # Instead, split source_text_raw on \n and normalise each line individually.
    # Only attempted when source_text_raw contains at least one newline.
    if source_text_raw and '\n' in source_text_raw:
        raw_lines = source_text_raw.split('\n')
        norm_lines = [normalise(l) for l in raw_lines]
        # Strip trailing empty lines (e.g. trailing newline in extraction)
        while norm_lines and not norm_lines[-1]:
            norm_lines.pop()
        n = len(norm_lines)
        if n > 1:
            paragraphs = shape.text_frame.paragraphs
            para_texts = [normalise(paragraph_text(p)) for p in paragraphs]
            for start in range(len(paragraphs) - n + 1):
                if para_texts[start:start + n] == norm_lines:
                    replace_paragraph_range(shape, start, n, translated_text)
                    log.info(
                        "%-24s seg=%-12s  slide=%s  shape=%s  paras=%d-%d",
                        f"{tier_prefix}T3b_RANGE", segment_id, slide_id,
                        label, start, start + n - 1,
                    )
                    return {
                        "result": f"{tier_prefix}T3b_PARA_RANGE",
                        "shape": label,
                        "para_start": start,
                        "para_end": start + n - 1,
                    }

    return None


# ---------------------------------------------------------------------------
# Speaker notes matching
# ---------------------------------------------------------------------------

def _match_notes(slide, source_text: str, translated_text: str,
                 segment_id: str) -> dict:
    """
    Match and replace text in a slide's speaker notes.
    Uses Tiers 1, 3, and 3b via the standard _match_shape helper.
    """
    try:
        notes_tf = slide.notes_slide.notes_text_frame
    except AttributeError:
        log.warning("NO NOTES  seg=%-12s — slide has no notes_slide.", segment_id)
        return {"segment_id": segment_id, "result": "NO_MATCH",
                "source_text": source_text}

    slide_id = getattr(slide, 'slide_id', '?')
    proxy = _NotesProxy(notes_tf, slide_id)
    norm_source = normalise(source_text)

    result = _match_shape(proxy, norm_source, translated_text,
                          segment_id, slide_id, tier_prefix="NOTES_",
                          source_text_raw=source_text)
    if result:
        result["segment_id"] = segment_id
        return result

    log.warning("NO MATCH (notes)  seg=%-12s  source_text=%r",
                segment_id, source_text[:80])
    return {"segment_id": segment_id, "result": "NO_MATCH",
            "source_text": source_text}


# ---------------------------------------------------------------------------
# Main find-and-replace orchestration
# ---------------------------------------------------------------------------

def find_and_replace(
    slide,
    source_text: str,
    translated_text: str,
    segment_id: str,
    state: dict,
) -> dict:
    norm_source = normalise(source_text)
    if not norm_source:
        return {"segment_id": segment_id, "result": "SKIP_EMPTY_SOURCE"}

    slide_id = getattr(slide, 'slide_id', '?')
    src_key = _source_key(norm_source)

    if src_key in state['matched_master_keys']:
        return {"segment_id": segment_id, "result": "MASTER_ALREADY_TRANSLATED"}
    if src_key in state['matched_layout_keys']:
        return {"segment_id": segment_id, "result": "LAYOUT_ALREADY_TRANSLATED"}

    # Tiers 1–3b: individual slide shapes
    for shape in _iter_shapes(slide.shapes):
        result = _match_shape(shape, norm_source, translated_text,
                              segment_id, slide_id,
                              source_text_raw=source_text)
        if result:
            result["segment_id"] = segment_id
            return result

    # Tier 4: slide layout shapes
    try:
        layout_shapes = slide.slide_layout.shapes
    except AttributeError:
        layout_shapes = []

    for shape in _iter_shapes(layout_shapes):
        if not shape.has_text_frame:
            continue
        para_ids = {id(p._p) for p in shape.text_frame.paragraphs}
        if para_ids & state['modified_layout_ids']:
            continue
        result = _match_shape(shape, norm_source, translated_text,
                              segment_id, slide_id, tier_prefix="LAYOUT_",
                              source_text_raw=source_text)
        if result:
            for p in shape.text_frame.paragraphs:
                state['modified_layout_ids'].add(id(p._p))
            state['matched_layout_keys'].add(src_key)
            result["segment_id"] = segment_id
            result["note"] = "Layout shape — applies to all slides using this layout."
            return result

    # Tier 5: slide master shapes
    try:
        master_shapes = slide.slide_layout.slide_master.shapes
    except AttributeError:
        master_shapes = []

    for shape in _iter_shapes(master_shapes):
        if not shape.has_text_frame:
            continue
        para_ids = {id(p._p) for p in shape.text_frame.paragraphs}
        if para_ids & state['modified_master_ids']:
            continue
        result = _match_shape(shape, norm_source, translated_text,
                              segment_id, slide_id, tier_prefix="MASTER_",
                              source_text_raw=source_text)
        if result:
            for p in shape.text_frame.paragraphs:
                state['modified_master_ids'].add(id(p._p))
            state['matched_master_keys'].add(src_key)
            result["segment_id"] = segment_id
            result["note"] = "Master shape — applies to all slides in the presentation."
            return result

    log.warning("NO MATCH  seg=%-12s  source_text=%r",
                segment_id, source_text[:80])
    return {"segment_id": segment_id, "result": "NO_MATCH",
            "source_text": source_text}


# ---------------------------------------------------------------------------
# HTML mop-up report
# ---------------------------------------------------------------------------

def _write_html_report(report: dict, output_path: Path,
                       source_filename: str) -> None:
    no_match = report.get("no_match_segments", [])
    summary  = report.get("summary", {})
    ts       = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")

    rows_html = ""
    for item in no_match:
        slide_num  = item.get("slide", "—")
        seg_id     = item.get("segment_id", "—")
        src        = item.get("source_text", "")
        translated = item.get("translated_text", "")
        rows_html += f"""
        <tr>
          <td class="slide">{html.escape(str(slide_num))}</td>
          <td class="seg">{html.escape(seg_id)}</td>
          <td class="src">{html.escape(src).replace(chr(10), '<br>')}</td>
          <td class="tgt">{html.escape(translated).replace(chr(10), '<br>')}</td>
        </tr>"""

    total      = summary.get("translated", 0)
    no_match_n = summary.get("NO_MATCH", 0)
    handled_n  = total - no_match_n

    page = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Manual Corrections — {html.escape(source_filename)}</title>
  <style>
    *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
      font-size: 14px; color: #1a1a2e; background: #f4f6fb; padding: 32px 24px;
    }}
    header {{ max-width: 1100px; margin: 0 auto 28px; }}
    h1 {{ font-size: 22px; font-weight: 700; margin-bottom: 6px; }}
    .meta {{ color: #555; font-size: 13px; margin-bottom: 18px; }}
    .stats {{ display: flex; gap: 16px; flex-wrap: wrap; margin-bottom: 28px; }}
    .stat {{
      background: #fff; border: 1px solid #dde3f0; border-radius: 8px;
      padding: 14px 20px; min-width: 140px;
    }}
    .stat .label {{ font-size: 11px; text-transform: uppercase;
                   letter-spacing: .06em; color: #777; margin-bottom: 4px; }}
    .stat .value {{ font-size: 26px; font-weight: 700; }}
    .stat.warn  .value {{ color: #c0392b; }}
    .stat.ok    .value {{ color: #27ae60; }}
    .notice {{
      background: #fff8e1; border-left: 4px solid #f39c12;
      border-radius: 4px; padding: 12px 16px; margin-bottom: 24px;
      max-width: 1100px; margin-left: auto; margin-right: auto;
      font-size: 13px; line-height: 1.6;
    }}
    .table-wrap {{ max-width: 1100px; margin: 0 auto; overflow-x: auto; }}
    table {{
      width: 100%; border-collapse: collapse; background: #fff;
      border-radius: 8px; overflow: hidden;
      box-shadow: 0 1px 4px rgba(0,0,0,.08);
    }}
    thead th {{
      background: #1a1a2e; color: #fff; font-size: 12px;
      text-transform: uppercase; letter-spacing: .06em;
      padding: 12px 14px; text-align: left;
    }}
    tbody tr:nth-child(odd)  {{ background: #fff; }}
    tbody tr:nth-child(even) {{ background: #f9fafc; }}
    tbody tr:hover {{ background: #eef2ff; }}
    td {{ padding: 10px 14px; vertical-align: top; line-height: 1.5;
         border-bottom: 1px solid #eee; }}
    td.slide {{ font-weight: 700; width: 60px; text-align: center; }}
    td.seg   {{ font-family: monospace; font-size: 12px; color: #555;
               white-space: nowrap; }}
    td.src   {{ color: #333; }}
    td.tgt   {{ color: #1a6b3a; font-weight: 500; }}
    .empty   {{ text-align: center; padding: 40px; color: #888; font-size: 16px; }}
  </style>
</head>
<body>
  <header>
    <h1>Manual Corrections Required</h1>
    <p class="meta">Source: {html.escape(source_filename)} &nbsp;·&nbsp; Generated: {ts}</p>
    <div class="stats">
      <div class="stat ok">
        <div class="label">Auto-applied</div>
        <div class="value">{handled_n}</div>
      </div>
      <div class="stat warn">
        <div class="label">Needs manual fix</div>
        <div class="value">{no_match_n}</div>
      </div>
      <div class="stat">
        <div class="label">Total segments</div>
        <div class="value">{total}</div>
      </div>
    </div>
    <div class="notice">
      <strong>How to use this file:</strong> Each row below is a text segment
      the pipeline could not apply automatically. The most common cause is text
      embedded in a raster image or SmartArt diagram. Open the translated PPTX,
      go to the slide shown, find the text in the <em>English (source)</em>
      column, and replace it with the <em>French (target)</em> text.
      Speaker notes are accessible via PowerPoint's Notes view (View &rarr; Notes).
    </div>
  </header>
  <div class="table-wrap">
    <table>
      <thead>
        <tr>
          <th>Slide</th>
          <th>Segment</th>
          <th>English (source — find this in the slide)</th>
          <th>French (target — replace with this)</th>
        </tr>
      </thead>
      <tbody>
        {"<tr><td colspan='4' class='empty'>No manual corrections required — all segments were applied automatically.</td></tr>" if not no_match else rows_html}
      </tbody>
    </table>
  </div>
</body>
</html>"""

    output_path.write_text(page, encoding="utf-8")
    log.info("Manual corrections report written to %s  (%d items)",
             output_path, no_match_n)


# ---------------------------------------------------------------------------
# Core reconstruction
# ---------------------------------------------------------------------------

def reconstruct(pptx_path: Path, json_path: Path, output_path: Path) -> dict:
    log.info("Loading presentation: %s", pptx_path)
    prs = Presentation(str(pptx_path))

    log.info("Loading translation JSON: %s", json_path)
    with json_path.open(encoding="utf-8") as fh:
        data = json.load(fh)

    segments = data.get("segments", [])
    log.info("Total segments in JSON: %d", len(segments))

    slide_map = {i + 1: slide for i, slide in enumerate(prs.slides)}
    log.info("Presentation has %d slides.", len(slide_map))

    state = {
        'modified_layout_ids': set(),
        'modified_master_ids': set(),
        'matched_layout_keys': set(),
        'matched_master_keys': set(),
    }

    results = []
    counts = {
        "total": 0, "translated": 0,
        "skipped_no_translation": 0, "skipped_do_not_translate": 0,
        "skipped_slide_out_of_range": 0,
        "skipped_not_text_box": 0,
        "T1_EXACT": 0, "T1b_FUZZY": 0, "T2_SUBSTRING": 0,
        "T3_PARAGRAPH": 0, "T3b_PARA_RANGE": 0,
        "LAYOUT_T1_EXACT": 0, "LAYOUT_T1b_FUZZY": 0,
        "LAYOUT_T2_SUBSTRING": 0, "LAYOUT_T3_PARAGRAPH": 0,
        "LAYOUT_T3b_PARA_RANGE": 0,
        "MASTER_T1_EXACT": 0, "MASTER_T1b_FUZZY": 0,
        "MASTER_T2_SUBSTRING": 0, "MASTER_T3_PARAGRAPH": 0,
        "MASTER_T3b_PARA_RANGE": 0,
        "LAYOUT_ALREADY_TRANSLATED": 0, "MASTER_ALREADY_TRANSLATED": 0,
        "NOTES_T1_EXACT": 0, "NOTES_T3_PARAGRAPH": 0,
        "NOTES_T3b_PARA_RANGE": 0,
        "NO_MATCH": 0, "SKIP_EMPTY_SOURCE": 0,
    }

    for seg in segments:
        counts["total"] += 1
        segment_id   = seg.get("segment_id", "UNKNOWN")
        slide_number = seg.get("slide_or_page")
        source_text  = seg.get("source_text", "")
        translated   = seg.get("translated_text")
        status       = seg.get("translation_status", "")
        element_type = seg.get("element_type", "")

        if not translated:
            counts["skipped_no_translation"] += 1
            continue
        if status == "DO_NOT_TRANSLATE":
            counts["skipped_do_not_translate"] += 1
            continue

        if seg.get("is_in_text_box") is False:
            if element_type not in _ALWAYS_PROCESS_TYPES:
                log.debug("SKIP_NOT_TEXT_BOX  seg=%-12s  type=%-16s  source=%r",
                          segment_id, element_type, source_text[:60])
                counts["skipped_not_text_box"] += 1
                results.append({
                    "segment_id": segment_id,
                    "result": "SKIP_NOT_TEXT_BOX",
                    "slide": slide_number,
                })
                continue
            log.debug("OVERRIDE_NOT_TEXT_BOX  seg=%-12s  type=%-16s  source=%r",
                      segment_id, element_type, source_text[:60])

        if slide_number not in slide_map:
            log.warning("Slide %s out of range for seg %s — skipping.",
                        slide_number, segment_id)
            counts["skipped_slide_out_of_range"] += 1
            results.append({"segment_id": segment_id,
                            "result": "SKIP_SLIDE_OUT_OF_RANGE",
                            "slide": slide_number})
            continue

        slide = slide_map[slide_number]
        counts["translated"] += 1

        if element_type == 'speaker_note':
            outcome = _match_notes(slide, source_text, translated, segment_id)
        else:
            outcome = find_and_replace(
                slide, source_text, translated, segment_id, state
            )

        outcome["slide"] = slide_number
        if outcome.get("result") == "NO_MATCH":
            outcome["translated_text"] = translated
        results.append(outcome)
        counts[outcome.get("result", "NO_MATCH")] = (
            counts.get(outcome.get("result", "NO_MATCH"), 0) + 1
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    log.info("Saving translated presentation to %s", output_path)
    prs.save(str(output_path))
    log.info("Save complete.")

    slide_matched  = (counts["T1_EXACT"] + counts["T1b_FUZZY"]
                      + counts["T2_SUBSTRING"] + counts["T3_PARAGRAPH"]
                      + counts["T3b_PARA_RANGE"])
    layout_matched = (counts["LAYOUT_T1_EXACT"] + counts["LAYOUT_T1b_FUZZY"]
                      + counts["LAYOUT_T2_SUBSTRING"] + counts["LAYOUT_T3_PARAGRAPH"]
                      + counts["LAYOUT_T3b_PARA_RANGE"])
    master_matched = (counts["MASTER_T1_EXACT"] + counts["MASTER_T1b_FUZZY"]
                      + counts["MASTER_T2_SUBSTRING"] + counts["MASTER_T3_PARAGRAPH"]
                      + counts["MASTER_T3b_PARA_RANGE"])
    notes_matched  = (counts["NOTES_T1_EXACT"] + counts["NOTES_T3_PARAGRAPH"]
                      + counts["NOTES_T3b_PARA_RANGE"])
    already_handled = (counts["LAYOUT_ALREADY_TRANSLATED"]
                       + counts["MASTER_ALREADY_TRANSLATED"])
    total_handled  = (slide_matched + layout_matched + master_matched
                      + notes_matched + already_handled)
    no_match       = counts["NO_MATCH"]
    handle_pct     = (total_handled / counts["translated"] * 100) if counts["translated"] else 0
    fail_rate      = no_match / counts["translated"] if counts["translated"] else 0

    log.info("=" * 60)
    log.info("RECONSTRUCTION SUMMARY")
    log.info("  Total segments              : %d", counts["total"])
    log.info("  Segments with translation   : %d", counts["translated"])
    log.info("  Slide-level matches")
    log.info("    Tier 1  (exact)           : %d", counts["T1_EXACT"])
    log.info("    Tier 1b (page-tolerant)   : %d", counts["T1b_FUZZY"])
    log.info("    Tier 2  (substring)       : %d", counts["T2_SUBSTRING"])
    log.info("    Tier 3  (paragraph)       : %d", counts["T3_PARAGRAPH"])
    log.info("    Tier 3b (para range)      : %d", counts["T3b_PARA_RANGE"])
    log.info("  Layout-level matches (Tier 4)")
    log.info("    T4 exact                  : %d", counts["LAYOUT_T1_EXACT"])
    log.info("    T4 page-tolerant          : %d", counts["LAYOUT_T1b_FUZZY"])
    log.info("    T4 substring              : %d", counts["LAYOUT_T2_SUBSTRING"])
    log.info("    T4 paragraph              : %d", counts["LAYOUT_T3_PARAGRAPH"])
    log.info("    T4 para range             : %d", counts["LAYOUT_T3b_PARA_RANGE"])
    log.info("  Master-level matches (Tier 5)")
    log.info("    T5 exact                  : %d", counts["MASTER_T1_EXACT"])
    log.info("    T5 page-tolerant          : %d", counts["MASTER_T1b_FUZZY"])
    log.info("    T5 substring              : %d", counts["MASTER_T2_SUBSTRING"])
    log.info("    T5 paragraph              : %d", counts["MASTER_T3_PARAGRAPH"])
    log.info("    T5 para range             : %d", counts["MASTER_T3b_PARA_RANGE"])
    log.info("  Notes-level matches")
    log.info("    Notes exact               : %d", counts["NOTES_T1_EXACT"])
    log.info("    Notes paragraph           : %d", counts["NOTES_T3_PARAGRAPH"])
    log.info("    Notes para range          : %d", counts["NOTES_T3b_PARA_RANGE"])
    log.info("  Already handled (deduped)")
    log.info("    Layout already translated : %d", counts["LAYOUT_ALREADY_TRANSLATED"])
    log.info("    Master already translated : %d", counts["MASTER_ALREADY_TRANSLATED"])
    log.info("  Total handled               : %d  (%.1f%%)", total_handled, handle_pct)
    log.info("  No match (manual required)  : %d  (%.1f%%)", no_match, fail_rate * 100)
    log.info("  Skipped (no translation)    : %d", counts["skipped_no_translation"])
    log.info("  Skipped (do not trans.)     : %d", counts["skipped_do_not_translate"])
    log.info("  Skipped (not text box)      : %d", counts["skipped_not_text_box"])
    log.info("  Skipped (out of range)      : %d", counts["skipped_slide_out_of_range"])
    log.info("=" * 60)

    return {
        "summary": counts,
        "handle_rate_pct": round(handle_pct, 2),
        "no_match_rate_pct": round(fail_rate * 100, 2),
        "results": results,
        "no_match_segments": [r for r in results if r.get("result") == "NO_MATCH"],
        "known_inaccessible_note": (
            "Persistent NO_MATCH segments are almost always text embedded in "
            "raster images or SmartArt, which python-pptx cannot access. "
            "See the manual corrections HTML report for actionable detail."
        ),
    }


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="CIPS Translation Pipeline — PPTX Reconstruction Script"
    )
    parser.add_argument("--source", required=True,
                        help="Path to the source PPTX file.")
    parser.add_argument("--translations", required=True,
                        help="Path to the Agent 3 translation JSON.")
    parser.add_argument("--qa", required=False, default=None,
                        help="Path to the Agent 4 QA JSON (optional).")
    parser.add_argument("--output", required=True,
                        help="Destination path for the translated PPTX.")
    parser.add_argument(
        "--failure-threshold", type=float, default=0.30,
        help=(
            "Exit with code 2 if the NO_MATCH rate exceeds this fraction. "
            "Default 0.30 (30%%). Increase for decks with significant "
            "embedded-image content that is structurally inaccessible."
        ),
    )
    args = parser.parse_args()

    source_path = Path(args.source)
    json_path   = Path(args.translations)
    output_path = Path(args.output)

    if not source_path.is_file():
        log.error("Source PPTX not found: %s", source_path)
        sys.exit(1)
    if not json_path.is_file():
        log.error("Translation JSON not found: %s", json_path)
        sys.exit(1)
    if args.qa:
        qa_path = Path(args.qa)
        if not qa_path.is_file():
            log.warning("QA file specified but not found: %s — proceeding without it.",
                        qa_path)

    report = reconstruct(source_path, json_path, output_path)

    report_path = (output_path.parent
                   / output_path.name.replace(".pptx", "_match_report.json"))
    report_path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    log.info("Match report written to %s", report_path)

    html_path = (output_path.parent
                 / output_path.name.replace(".pptx", "_manual_corrections.html"))
    _write_html_report(report, html_path, source_path.name)

    no_match_rate = (report["summary"].get("NO_MATCH", 0)
                     / max(report["summary"].get("translated", 1), 1))
    if no_match_rate > args.failure_threshold:
        log.error(
            "NO_MATCH rate %.1f%% exceeds threshold %.0f%% — "
            "review %s for required manual corrections.",
            no_match_rate * 100,
            args.failure_threshold * 100,
            html_path,
        )
        sys.exit(2)

    log.info("Reconstruction complete. Output: %s", output_path)


if __name__ == "__main__":
    main()
