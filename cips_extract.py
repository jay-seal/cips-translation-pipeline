"""
cips_extract.py
===============
CIPS Translation Pipeline — Source Document Extraction Script

Replaces Agent 1 in the automated workflow. Extracts all text from a PPTX
or DOCX source file using python-pptx / python-docx and outputs one
Agent 1-format JSON file per batch. Because extraction happens directly in
python-pptx, the source_text values are guaranteed to match exactly what
cips_reconstruct.py reads from the same file — eliminating the smart-quote
mismatches and hallucination risk that occur when an LLM agent processes a
binary file via URL.

Usage:
    python cips_extract.py \\
        --source  inputs/source.pptx \\
        --output-dir  inputs/ \\
        --locale  FR-FR \\
        --document-type  pptx \\
        --batch-size  20

Outputs per batch:
    inputs/agent1_slides_001_020.json
    inputs/agent1_slides_021_040.json
    ...

Also writes:
    inputs/extraction_manifest.json   — list of batch filenames and metadata

Exit codes:
    0   Success.
    1   Missing or invalid input.
"""

import argparse
import json
import logging
import re
import sys
import unicodedata
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from docx import Document
from docx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("cips_extract")

# ---------------------------------------------------------------------------
# Placeholder type constants (python-pptx PP_PLACEHOLDER enum values)
# ---------------------------------------------------------------------------
_PH_TITLE_TYPES    = frozenset([PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE])
_PH_SUBTITLE_TYPES = frozenset([PP_PLACEHOLDER.SUBTITLE])
_PH_BODY_TYPES     = frozenset([PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT])
_PH_FOOTER_TYPES   = frozenset([
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.SLIDE_NUMBER,
])

# Element types that are always editable — is_in_text_box = True
_ALWAYS_TEXT_BOX = frozenset([
    'slide_title', 'slide_subtitle', 'body_text', 'bullet_point',
    'heading', 'table_cell', 'text_box', 'caption', 'speaker_note',
])


# ---------------------------------------------------------------------------
# Shape iteration (mirrored from cips_reconstruct.py)
# ---------------------------------------------------------------------------

def _iter_shapes(shape_collection):
    for shape in shape_collection:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_group(shape)
        elif st == MSO_SHAPE_TYPE.TABLE:
            yield from _iter_table_cells(shape)
        else:
            yield ('shape', shape)


def _iter_group(group_shape):
    for shape in group_shape.shapes:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_group(shape)
        elif st == MSO_SHAPE_TYPE.TABLE:
            yield from _iter_table_cells(shape)
        else:
            yield ('shape', shape)


def _iter_table_cells(table_shape):
    for ri, row in enumerate(table_shape.table.rows):
        for ci, cell in enumerate(row.cells):
            yield ('table_cell', table_shape, ri, ci, cell)


# ---------------------------------------------------------------------------
# Element type determination
# ---------------------------------------------------------------------------

def _pptx_element_type(shape) -> str:
    """Derive element_type from a PPTX shape's placeholder format."""
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in _PH_TITLE_TYPES:
            return 'slide_title'
        if ph_type in _PH_SUBTITLE_TYPES:
            return 'slide_subtitle'
        if ph_type in _PH_FOOTER_TYPES:
            return 'footer'
        if ph_type in _PH_BODY_TYPES:
            return 'body_text'   # Refined to bullet_point at paragraph level
        return 'text_box'

    # Non-placeholder text boxes: short single-paragraph shapes are labels
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
        if len(text) <= 60 and len(paras) <= 1:
            return 'label'
    return 'text_box'


def _para_element_type(base_type: str, para) -> str:
    """Refine element_type at paragraph level for body placeholders."""
    if base_type != 'body_text':
        return base_type
    # Indented paragraphs or paragraphs with bullet XML are bullet_points
    if para.level > 0:
        return 'bullet_point'
    pPr = para._p.find(qn('a:pPr'))
    if pPr is not None and pPr.get('indent') is not None:
        return 'bullet_point'
    return 'body_text'


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------

def _full_para_text(para) -> str:
    """Return full paragraph text including field elements."""
    return ''.join(
        node.text or ''
        for node in para._p.iter()
        if node.tag in (qn('a:t'), qn('a:delText'))
    )


def extract_pptx(path: Path, locale: str) -> list[dict]:
    """
    Extract all segments from a PPTX file.
    Returns a flat list of segment dicts sorted by slide number.
    """
    prs = Presentation(str(path))
    source_filename = path.name
    total_slides = len(prs.slides)
    segments = []
    seg_counter = [0]   # mutable so nested helpers can increment

    def _make_seg(slide_num, elem_type, elem_id, source_text, layout_ctx,
                  is_table=False):
        text = source_text.strip()
        if not text:
            return
        seg_counter[0] += 1
        is_in_tb = elem_type in _ALWAYS_TEXT_BOX
        segments.append({
            "segment_id":         f"SEG-{seg_counter[0]:03d}",
            "slide_or_page":      slide_num,
            "element_type":       elem_type,
            "element_id":         elem_id,
            "source_text":        text,
            "character_count":    len(text),
            "is_in_text_box":     is_in_tb,
            "is_in_table":        is_table,
            "layout_context":     layout_ctx,
            "translation_status": "PENDING",
            "translated_text":    None,
            "flag_status":        None,
            "flag_options":       None,
            "expansion_risk":     None,
        })

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        shape_counter = 0

        for item in _iter_shapes(slide.shapes):
            if item[0] == 'table_cell':
                _, table_shape, ri, ci, cell = item
                if not cell.text_frame.text.strip():
                    continue
                cell_text = '\n'.join(
                    p.text.strip()
                    for p in cell.text_frame.paragraphs
                    if p.text.strip()
                )
                elem_id = (f"slide_{slide_num}_{table_shape.name}"
                           f"_row{ri}_col{ci}").replace(' ', '_')
                _make_seg(slide_num, 'table_cell', elem_id, cell_text,
                          f"Table row {ri+1}, column {ci+1}", is_table=True)

            else:  # regular shape
                _, shape = item
                if not shape.has_text_frame:
                    continue
                base_type = _pptx_element_type(shape)
                shape_counter += 1

                # Treat each paragraph as a separate segment for body/bullet
                # shapes; collapse single-para shapes to one segment.
                paras = [p for p in shape.text_frame.paragraphs
                         if _full_para_text(p).strip()]
                if not paras:
                    continue

                shape_label = shape.name.replace(' ', '_')
                is_multi_para_body = (base_type in ('body_text',)
                                      and len(paras) > 1)

                if is_multi_para_body:
                    for p_idx, para in enumerate(paras):
                        text = _full_para_text(para).strip()
                        if not text:
                            continue
                        etype = _para_element_type(base_type, para)
                        elem_id = (f"slide_{slide_num}_{shape_label}"
                                   f"_para{p_idx+1}")
                        _make_seg(slide_num, etype, elem_id, text,
                                  f"Paragraph {p_idx+1} in {shape.name}")
                else:
                    # Single-paragraph or non-body shape: one segment
                    text_parts = [_full_para_text(p) for p in paras]
                    full_text = '\n'.join(t.strip() for t in text_parts
                                         if t.strip())
                    if not full_text:
                        continue
                    etype = base_type
                    if len(paras) == 1:
                        etype = _para_element_type(base_type, paras[0])
                    elem_id = f"slide_{slide_num}_{shape_label}"
                    _make_seg(slide_num, etype, elem_id, full_text,
                              f"Shape: {shape.name}")

        # Speaker notes
        try:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = '\n'.join(
                p.text.strip() for p in notes_tf.paragraphs
                if p.text.strip()
            )
            if notes_text:
                _make_seg(slide_num, 'speaker_note',
                          f"slide_{slide_num}_notes",
                          notes_text, "Speaker notes")
        except Exception:
            pass

    log.info("PPTX extraction complete: %d segments from %d slides",
             len(segments), total_slides)
    return segments, total_slides, source_filename


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------

def _docx_para_text(para) -> str:
    """Full paragraph text including field text."""
    return ''.join(
        node.text or ''
        for node in para._p.iter()
        if node.tag in (qn('w:t'), qn('w:delText'))
    )


def _docx_element_type(para) -> str:
    style_name = (para.style.name or '').lower()
    if 'heading' in style_name or 'title' in style_name:
        return 'heading'
    if 'list' in style_name or 'bullet' in style_name:
        return 'bullet_point'
    return 'body_text'


def extract_docx(path: Path, locale: str) -> tuple:
    """
    Extract all segments from a DOCX file.
    Returns (segments, approx_page_count, source_filename).
    Pages are approximated at ~400 words per page.
    """
    doc = Document(str(path))
    source_filename = path.name
    segments = []
    seg_counter = [0]
    para_counter = [0]

    def _make_seg(page_num, elem_type, elem_id, text, ctx,
                  is_in_tb=True, is_table=False):
        text = text.strip()
        if not text:
            return
        seg_counter[0] += 1
        segments.append({
            "segment_id":         f"SEG-{seg_counter[0]:03d}",
            "slide_or_page":      page_num,
            "element_type":       elem_type,
            "element_id":         elem_id,
            "source_text":        text,
            "character_count":    len(text),
            "is_in_text_box":     is_in_tb,
            "is_in_table":        is_table,
            "layout_context":     ctx,
            "translation_status": "PENDING",
            "translated_text":    None,
            "flag_status":        None,
            "flag_options":       None,
            "expansion_risk":     None,
        })

    # Estimate cumulative word count for page approximation
    cumulative_words = 0
    WORDS_PER_PAGE = 350

    for i, para in enumerate(doc.paragraphs):
        text = _docx_para_text(para).strip()
        if not text:
            continue
        para_counter[0] += 1
        cumulative_words += len(text.split())
        page_num = max(1, (cumulative_words // WORDS_PER_PAGE) + 1)
        etype = _docx_element_type(para)
        _make_seg(page_num, etype, f"page_{page_num}_para_{para_counter[0]}",
                  text, "Document paragraph")

    # Tables
    for t_idx, table in enumerate(doc.tables):
        for ri, row in enumerate(table.rows):
            for ci, cell in enumerate(row.cells):
                cell_text = '\n'.join(
                    _docx_para_text(p).strip()
                    for p in cell.paragraphs
                    if _docx_para_text(p).strip()
                )
                if not cell_text:
                    continue
                _make_seg(1, 'table_cell',
                          f"table_{t_idx+1}_row{ri}_col{ci}",
                          cell_text, f"Table {t_idx+1}, row {ri+1}, col {ci+1}",
                          is_table=True)

    # Footers
    for section in doc.sections:
        for ftr_para in section.footer.paragraphs:
            text = _docx_para_text(ftr_para).strip()
            if text:
                _make_seg(1, 'footer', 'footer', text, "Footer",
                          is_in_tb=False)
                break  # One footer segment per document

    total_pages = max(1, (cumulative_words // WORDS_PER_PAGE) + 1)
    log.info("DOCX extraction complete: %d segments (~%d pages)",
             len(segments), total_pages)
    return segments, total_pages, source_filename


# ---------------------------------------------------------------------------
# Batch writing
# ---------------------------------------------------------------------------

def write_batches(
    segments: list,
    total_pages: int,
    source_filename: str,
    doc_type: str,
    locale: str,
    batch_size: int,
    output_dir: Path,
) -> list[str]:
    """
    Split segments into batches by slide/page and write one JSON per batch.
    Returns list of output filenames (not full paths).
    """
    # Group segments by page number
    pages = sorted({s['slide_or_page'] for s in segments})
    all_pages = list(range(1, total_pages + 1))

    batch_filenames = []
    batch_num = 0
    page_iter = iter(range(1, total_pages + 1, batch_size))

    for batch_start in range(1, total_pages + 1, batch_size):
        batch_end = min(batch_start + batch_size - 1, total_pages)
        batch_segs = [s for s in segments
                      if batch_start <= s['slide_or_page'] <= batch_end]

        # Renumber segment IDs within this batch
        for i, seg in enumerate(batch_segs, start=1):
            seg['segment_id'] = f"SEG-{i:03d}"

        range_str = f"{batch_start:03d}_{batch_end:03d}"
        if doc_type == 'docx':
            filename = f"agent1_pages_{range_str}.json"
        else:
            filename = f"agent1_slides_{range_str}.json"

        payload = {
            "document_metadata": {
                "source_filename":      source_filename,
                "document_type":        doc_type,
                "target_locale":        locale,
                "total_slides_or_pages": total_pages,
                "batch_slide_range":    f"{batch_start}-{batch_end}",
                "extraction_timestamp": datetime.now(timezone.utc).isoformat(),
            },
            "segments": batch_segs,
        }

        out_path = output_dir / filename
        out_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding='utf-8',
        )
        batch_filenames.append(filename)
        log.info("Wrote %s  (%d segments)", filename, len(batch_segs))

    return batch_filenames


def write_manifest(
    batch_filenames: list,
    source_filename: str,
    doc_type: str,
    locale: str,
    total_pages: int,
    batch_size: int,
    output_dir: Path,
    source_url: str = "",
) -> None:
    manifest = {
        "source_filename": source_filename,
        "source_url":      source_url or "",
        "document_type":   doc_type,
        "locale_code":     locale,
        "total_pages":     total_pages,
        "batch_size":      batch_size,
        "batch_count":     len(batch_filenames),
        "batch_files":     batch_filenames,
        "created_at":      datetime.now(timezone.utc).isoformat(),
    }
    manifest_path = output_dir / "extraction_manifest.json"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding='utf-8',
    )
    log.info("Manifest written: %s  (%d batches)", manifest_path,
             len(batch_filenames))


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="CIPS Translation Pipeline — Source Document Extraction"
    )
    parser.add_argument('--source', required=True,
                        help="Path to source PPTX or DOCX file.")
    parser.add_argument('--output-dir', required=True,
                        help="Directory to write batch JSON files.")
    parser.add_argument('--locale', required=True,
                        help="Target locale code, e.g. FR-FR.")
    parser.add_argument('--document-type', required=True,
                        choices=['pptx', 'docx'],
                        help="Source document format.")
    parser.add_argument('--batch-size', type=int, default=20,
                        help="Slides/pages per batch (default 20).")
    parser.add_argument('--source-url', default=None,
                        help="Original R2 URL of the source file. Stored in the manifest "
                             "so downstream workflows can retrieve it without it being "
                             "passed through every trigger.")
    parser.add_argument('--source-filename', default=None,
                        help="Original filename to embed in JSON metadata. "
                             "Overrides the name derived from --source path. "
                             "Use when the source file was downloaded with a "
                             "generic name (e.g. source.pptx).")
    args = parser.parse_args()

    source_path = Path(args.source)
    output_dir  = Path(args.output_dir)

    if not source_path.is_file():
        log.error("Source file not found: %s", source_path)
        sys.exit(1)
    output_dir.mkdir(parents=True, exist_ok=True)

    if args.document_type == 'pptx':
        segments, total_pages, source_filename = extract_pptx(
            source_path, args.locale)
    else:
        segments, total_pages, source_filename = extract_docx(
            source_path, args.locale)

    # Override the filename with the original name if provided
    if args.source_filename:
        source_filename = args.source_filename

    batch_filenames = write_batches(
        segments, total_pages, source_filename,
        args.document_type, args.locale, args.batch_size, output_dir,
    )
    write_manifest(
        batch_filenames, source_filename, args.document_type,
        args.locale, total_pages, args.batch_size, output_dir,
        source_url=args.source_url,
    )

    log.info("=" * 60)
    log.info("EXTRACTION COMPLETE")
    log.info("  Source         : %s", source_path)
    log.info("  Document type  : %s", args.document_type)
    log.info("  Total pages    : %d", total_pages)
    log.info("  Batch size     : %d", args.batch_size)
    log.info("  Batches written: %d", len(batch_filenames))
    log.info("  Output dir     : %s", output_dir)
    log.info("=" * 60)


if __name__ == '__main__':
    main()
